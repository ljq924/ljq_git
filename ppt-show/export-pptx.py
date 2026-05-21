from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DEEP_NAVY = RGBColor(0x0F, 0x1B, 0x3D)
DARK_VOID = RGBColor(0x0A, 0x0E, 0x27)
NEON_CYAN = RGBColor(0x5E, 0xDC, 0xF4)
NEON_PINK = RGBColor(0xF0, 0xA6, 0xCA)
NEON_YELLOW = RGBColor(0xF4, 0xD0, 0x3F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_LAVENDER = RGBColor(0xE2, 0xD5, 0xF2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_name='Tektur', font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

# ===== SLIDE 1: HERO =====
slide1 = prs.slides.add_slide(blank_layout)
add_bg(slide1, DARK_VOID)
add_rect(slide1, 0, 0, 13.333, 7.5, DARK_VOID)
add_textbox(slide1, 2, 0.8, 9, 0.6, 'COMPUTER SCIENCE STUDENT', 'Space Mono', 14, NEON_PINK, False, PP_ALIGN.CENTER)
add_textbox(slide1, 2, 1.8, 9, 2.5, 'HELLO\nWORLD', 'Tektur', 72, NEON_CYAN, True, PP_ALIGN.CENTER)
add_textbox(slide1, 3, 4.8, 7, 1.2, '热爱编程，探索技术的无限可能。\n从算法到全栈，用代码构建未来。', 'Chakra Petch', 18, RGBColor(0xAA, 0xBB, 0xCC), False, PP_ALIGN.CENTER)
add_textbox(slide1, 2.5, 6.2, 2.5, 0.5, 'FULL-STACK', 'Space Mono', 11, NEON_YELLOW, True, PP_ALIGN.CENTER)
add_textbox(slide1, 5.5, 6.2, 2.5, 0.5, 'OPEN SOURCE', 'Space Mono', 11, NEON_YELLOW, True, PP_ALIGN.CENTER)
add_textbox(slide1, 8.5, 6.2, 2.5, 0.5, 'AI EXPLORER', 'Space Mono', 11, NEON_YELLOW, True, PP_ALIGN.CENTER)

# ===== SLIDE 2: ABOUT ME =====
slide2 = prs.slides.add_slide(blank_layout)
add_bg(slide2, NEON_PINK)
add_textbox(slide2, 0.8, 0.5, 3, 0.5, 'PROFILE', 'Space Mono', 12, NEON_YELLOW, True, PP_ALIGN.LEFT)
add_textbox(slide2, 0.8, 1.2, 5, 1, '关于我', 'Tektur', 40, DEEP_NAVY, True, PP_ALIGN.LEFT)
add_textbox(slide2, 0.8, 2.5, 5, 3, '一名计算机科学专业在读学生，对软件开发和人工智能充满热情。喜欢用技术解决实际问题，相信代码可以改变世界。\n\n在课余时间积极参与开源项目和技术社区，不断学习新技术，追求成为全栈开发工程师。', 'Chakra Petch', 16, RGBColor(0x33, 0x33, 0x55), False, PP_ALIGN.LEFT)
pixel_box = add_rect(slide2, 7.5, 1.5, 4.5, 4.5, DEEP_NAVY, NEON_YELLOW, Pt(4))
add_textbox(slide2, 8.2, 2.5, 3, 1, '👾', 'Tektur', 80, NEON_CYAN, True, PP_ALIGN.CENTER)
add_textbox(slide2, 8.2, 4.2, 3, 0.8, 'CS Student', 'Space Mono', 14, NEON_PINK, True, PP_ALIGN.CENTER)

# ===== SLIDE 3: SKILLS =====
slide3 = prs.slides.add_slide(blank_layout)
add_bg(slide3, NEON_CYAN)
add_textbox(slide3, 4, 0.3, 5, 0.5, 'TECH STACK', 'Space Mono', 12, DEEP_NAVY, True, PP_ALIGN.CENTER)
add_textbox(slide3, 4, 0.9, 5, 1, '技术栈', 'Tektur', 40, DEEP_NAVY, True, PP_ALIGN.CENTER)

skills = [
    ('前端开发', 'React / Vue / TypeScript\nHTML5 / CSS3 / Tailwind'),
    ('后端开发', 'Node.js / Python / Java\nSpring Boot / Express'),
    ('AI / ML', 'PyTorch / TensorFlow\nNLP / CV / LLM'),
    ('工具 & DevOps', 'Git / Docker / Linux\nCI/CD / PostgreSQL'),
]
for i, (title, desc) in enumerate(skills):
    col = i % 4
    x = 0.8 + col * 3.1
    y = 2.2
    card = add_rect(slide3, x, y, 2.8, 4, WHITE, RGBColor(0xCC, 0xCC, 0xCC), Pt(1))
    add_textbox(slide3, x + 0.2, y + 0.3, 2.4, 0.6, title, 'Tektur', 18, DEEP_NAVY, True, PP_ALIGN.CENTER)
    add_textbox(slide3, x + 0.2, y + 1.2, 2.4, 2, desc, 'Chakra Petch', 12, RGBColor(0x55, 0x55, 0x77), False, PP_ALIGN.CENTER)

# ===== SLIDE 4: SKILL CHART =====
slide4 = prs.slides.add_slide(blank_layout)
add_bg(slide4, DARK_VOID)
add_textbox(slide4, 4, 0.3, 5, 0.5, 'SKILL LEVEL', 'Space Mono', 12, NEON_YELLOW, True, PP_ALIGN.CENTER)
add_textbox(slide4, 4, 0.9, 5, 1, '技能熟练度', 'Tektur', 40, NEON_CYAN, True, PP_ALIGN.CENTER)

chart_data = [
    ('Python', 90, NEON_CYAN),
    ('JavaScript', 85, NEON_PINK),
    ('React', 78, NEON_YELLOW),
    ('Java', 72, NEON_CYAN),
    ('Docker', 65, NEON_PINK),
]
bar_base_y = 6.5
bar_max_h = 4.0
for i, (name, val, color) in enumerate(chart_data):
    x = 1.5 + i * 2.2
    bar_h = bar_max_h * val / 100
    bar_y = bar_base_y - bar_h
    add_rect(slide4, x, bar_y, 1.5, bar_h, color)
    add_textbox(slide4, x - 0.2, bar_y - 0.5, 1.9, 0.5, str(val), 'Space Mono', 14, NEON_YELLOW, True, PP_ALIGN.CENTER)
    add_textbox(slide4, x - 0.3, bar_base_y + 0.1, 2.1, 0.5, name, 'Space Mono', 11, RGBColor(0x88, 0x99, 0xAA), False, PP_ALIGN.CENTER)

# ===== SLIDE 5: TIMELINE =====
slide5 = prs.slides.add_slide(blank_layout)
add_bg(slide5, NEON_PINK)
add_textbox(slide5, 4.5, 0.3, 4, 0.5, 'JOURNEY', 'Space Mono', 12, DEEP_NAVY, True, PP_ALIGN.CENTER)
add_textbox(slide5, 4.5, 0.9, 4, 1, '成长历程', 'Tektur', 40, DEEP_NAVY, True, PP_ALIGN.CENTER)

timeline = [
    ('2022', '踏入编程世界', '学习第一门编程语言 C 语言，打开计算机科学的大门。'),
    ('2023', '全栈探索', '深入学习 Web 开发，掌握前后端技术栈，完成首个完整项目。'),
    ('2024', 'AI 深入研究', '接触机器学习与深度学习，参与 AI 相关竞赛与科研项目。'),
    ('2025', '开源 & 实习', '贡献开源项目，积累工程经验，向专业开发者迈进。'),
]
for i, (year, title, desc) in enumerate(timeline):
    y = 2.0 + i * 1.3
    add_rect(slide5, 1, y, 1.2, 0.4, DEEP_NAVY)
    add_textbox(slide5, 1, y, 1.2, 0.4, year, 'Space Mono', 11, NEON_CYAN, True, PP_ALIGN.CENTER)
    add_textbox(slide5, 2.5, y - 0.1, 4, 0.5, title, 'Tektur', 18, DEEP_NAVY, True, PP_ALIGN.LEFT)
    add_textbox(slide5, 2.5, y + 0.4, 9, 0.6, desc, 'Chakra Petch', 13, RGBColor(0x33, 0x33, 0x55), False, PP_ALIGN.LEFT)

# ===== SLIDE 6: STATS =====
slide6 = prs.slides.add_slide(blank_layout)
add_bg(slide6, DARK_VOID)
add_textbox(slide6, 4.5, 0.3, 4, 0.5, 'ACHIEVEMENTS', 'Space Mono', 12, NEON_YELLOW, True, PP_ALIGN.CENTER)
add_textbox(slide6, 4.5, 0.9, 4, 1, '数据一览', 'Tektur', 40, NEON_CYAN, True, PP_ALIGN.CENTER)

stats = [
    ('15', '完成项目'),
    ('500+', 'GitHub Commits'),
    ('3', '竞赛获奖'),
    ('8', '掌握语言'),
]
for i, (num, label) in enumerate(stats):
    x = 1 + i * 3
    add_rect(slide6, x, 2.5, 2.5, 3.5, RGBColor(0x12, 0x22, 0x44), NEON_CYAN, Pt(2))
    add_textbox(slide6, x, 3.0, 2.5, 1.5, num, 'Tektur', 48, NEON_CYAN, True, PP_ALIGN.CENTER)
    add_textbox(slide6, x, 4.5, 2.5, 0.5, label, 'Space Mono', 12, NEON_PINK, True, PP_ALIGN.CENTER)

# ===== SLIDE 7: QUOTE =====
slide7 = prs.slides.add_slide(blank_layout)
add_bg(slide7, NEON_CYAN)
add_textbox(slide7, 2, 1.5, 9, 3, '" Talk is cheap.\nShow me the code. "\n\n代码是程序员最好的语言，每一行都在诉说对世界的理解与重塑。', 'Chakra Petch', 24, DEEP_NAVY, False, PP_ALIGN.CENTER)
add_rect(slide7, 5.5, 5.2, 2, 0.08, NEON_YELLOW)
add_textbox(slide7, 3, 5.5, 7, 0.5, '— Linus Torvalds', 'Space Mono', 14, RGBColor(0x55, 0x55, 0x77), False, PP_ALIGN.CENTER)

# ===== SLIDE 8: PROJECTS =====
slide8 = prs.slides.add_slide(blank_layout)
add_bg(slide8, SOFT_LAVENDER)
add_textbox(slide8, 4.5, 0.3, 4, 0.5, 'PROJECTS', 'Space Mono', 12, NEON_PINK, True, PP_ALIGN.CENTER)
add_textbox(slide8, 4.5, 0.9, 4, 1, '项目展示', 'Tektur', 40, DEEP_NAVY, True, PP_ALIGN.CENTER)

projects = [
    ('智能聊天助手', 'Python / LLM / FastAPI', '基于大语言模型的智能对话系统，支持多轮对话与知识库检索。', 'AI', False),
    ('全栈电商平台', 'React / Node.js / PostgreSQL', '完整的电商系统，包含用户认证、商品管理、支付集成和订单追踪。', 'Full-Stack', True),
    ('数据可视化面板', 'Vue / D3.js / ECharts', '实时数据监控仪表盘，支持自定义图表和动态数据源接入。', 'Data Viz', False),
]
for i, (name, tech, desc, tag, featured) in enumerate(projects):
    x = 0.8 + i * 4.1
    bg_color = DEEP_NAVY if featured else WHITE
    txt_color = WHITE if featured else DEEP_NAVY
    card = add_rect(slide8, x, 2.2, 3.8, 4.5, bg_color, RGBColor(0xCC, 0xCC, 0xCC), Pt(1))
    add_textbox(slide8, x + 0.3, 2.6, 3.2, 0.6, name, 'Tektur', 20, NEON_YELLOW if featured else DEEP_NAVY, True, PP_ALIGN.CENTER)
    add_textbox(slide8, x + 0.3, 3.3, 3.2, 0.5, tech, 'Space Mono', 10, NEON_CYAN if featured else RGBColor(0x88, 0x88, 0x88), False, PP_ALIGN.CENTER)
    add_textbox(slide8, x + 0.3, 4.0, 3.2, 2, desc, 'Chakra Petch', 12, RGBColor(0xAA, 0xAA, 0xCC) if featured else RGBColor(0x66, 0x66, 0x66), False, PP_ALIGN.CENTER)
    add_textbox(slide8, x + 0.3, 5.8, 3.2, 0.5, tag, 'Space Mono', 10, NEON_PINK if featured else DEEP_NAVY, True, PP_ALIGN.CENTER)

# ===== SLIDE 9: CTA =====
slide9 = prs.slides.add_slide(blank_layout)
add_bg(slide9, DARK_VOID)
add_textbox(slide9, 2, 1.5, 9, 2, "LET'S\nCONNECT", 'Tektur', 64, NEON_CYAN, True, PP_ALIGN.CENTER)
add_textbox(slide9, 3, 3.8, 7, 1, '期待与志同道合的伙伴交流合作，一起用代码创造更多可能。', 'Chakra Petch', 18, RGBColor(0xAA, 0xBB, 0xCC), False, PP_ALIGN.CENTER)
add_textbox(slide9, 3, 5.2, 3, 0.5, 'GitHub', 'Tektur', 16, NEON_CYAN, True, PP_ALIGN.CENTER)
add_textbox(slide9, 7, 5.2, 3, 0.5, 'Email Me', 'Tektur', 16, NEON_PINK, True, PP_ALIGN.CENTER)
add_textbox(slide9, 3, 6.2, 2, 0.4, 'LinkedIn', 'Space Mono', 11, NEON_PINK, False, PP_ALIGN.CENTER)
add_textbox(slide9, 5.5, 6.2, 2, 0.4, 'Blog', 'Space Mono', 11, NEON_PINK, False, PP_ALIGN.CENTER)
add_textbox(slide9, 8, 6.2, 2, 0.4, 'Resume', 'Space Mono', 11, NEON_PINK, False, PP_ALIGN.CENTER)

output_path = r'f:\zuomian\项目\ppt-show\personal-intro.pptx'
prs.save(output_path)
print(f'PPTX saved to: {output_path}')
